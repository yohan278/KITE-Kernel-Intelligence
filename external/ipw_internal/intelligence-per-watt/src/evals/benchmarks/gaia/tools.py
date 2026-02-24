# benchmarks/gaia/tools.py
"""
Tools for handling various file types in GAIA benchmark.

Provides functions to read and process files commonly used in GAIA questions:
- Excel (xlsx, xls)
- Images (png, jpeg, jpg)
- Audio (mp3, wav)
- Documents (pdf, docx, pptx)
- Data formats (csv, json, jsonld, xml)
- Code (py, zip)
- Other (pdb, txt)
"""
from pathlib import Path
from typing import Any, Dict, List, Optional
import logging

logger = logging.getLogger(__name__)


def read_text_file(file_path: str) -> str:
    """Read a text file and return its contents.
    
    Args:
        file_path: Path to the text file
        
    Returns:
        File contents as string
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        return f"Error reading text file: {e}"


def read_csv_file(file_path: str) -> str:
    """Read a CSV file and return its contents as a formatted string.
    
    Args:
        file_path: Path to the CSV file
        
    Returns:
        CSV contents as formatted string with row information
    """
    try:
        import pandas as pd
        df = pd.read_csv(file_path)
        result = f"CSV file with {len(df)} rows and {len(df.columns)} columns\n"
        result += f"Columns: {', '.join(df.columns.tolist())}\n\n"
        result += f"First 10 rows:\n{df.head(10).to_string()}\n\n"
        result += f"Summary statistics:\n{df.describe().to_string()}"
        return result
    except Exception as e:
        return f"Error reading CSV file: {e}"


def read_excel_file(file_path: str, sheet_name: Optional[str] = None) -> str:
    """Read an Excel file and return its contents.
    
    Args:
        file_path: Path to the Excel file
        sheet_name: Specific sheet to read (None for first sheet)
        
    Returns:
        Excel contents as formatted string
    """
    try:
        import pandas as pd
        
        # Get all sheet names first
        excel_file = pd.ExcelFile(file_path)
        sheet_names = excel_file.sheet_names
        
        result = f"Excel file with {len(sheet_names)} sheet(s): {', '.join(sheet_names)}\n\n"
        
        # Read the specified sheet or first sheet
        if sheet_name is None:
            sheet_name = sheet_names[0]
        
        df = pd.read_excel(file_path, sheet_name=sheet_name)
        result += f"Reading sheet: {sheet_name}\n"
        result += f"{len(df)} rows and {len(df.columns)} columns\n"
        result += f"Columns: {', '.join(df.columns.tolist())}\n\n"
        result += f"First 10 rows:\n{df.head(10).to_string()}\n\n"
        result += f"Summary statistics:\n{df.describe().to_string()}"
        
        return result
    except Exception as e:
        return f"Error reading Excel file: {e}"


def read_json_file(file_path: str) -> str:
    """Read a JSON or JSON-LD file and return its contents.
    
    Args:
        file_path: Path to the JSON file
        
    Returns:
        JSON contents as formatted string
    """
    try:
        import json
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # Pretty print the JSON
        json_str = json.dumps(data, indent=2, ensure_ascii=False)
        
        result = f"JSON file contents:\n{json_str}\n\n"
        result += f"Type: {type(data).__name__}\n"
        
        if isinstance(data, dict):
            result += f"Keys: {', '.join(data.keys())}"
        elif isinstance(data, list):
            result += f"Length: {len(data)} items"
        
        return result
    except Exception as e:
        return f"Error reading JSON file: {e}"


def read_pdf_file(file_path: str, max_pages: int = 10) -> str:
    """Read a PDF file and extract text.
    
    Args:
        file_path: Path to the PDF file
        max_pages: Maximum number of pages to read
        
    Returns:
        Extracted text from PDF
    """
    try:
        import PyPDF2
        
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            num_pages = len(pdf_reader.pages)
            
            result = f"PDF file with {num_pages} pages\n\n"
            
            pages_to_read = min(num_pages, max_pages)
            for i in range(pages_to_read):
                page = pdf_reader.pages[i]
                text = page.extract_text()
                result += f"--- Page {i+1} ---\n{text}\n\n"
            
            if num_pages > max_pages:
                result += f"(Showing first {max_pages} of {num_pages} pages)"
            
            return result
    except ImportError:
        return "PyPDF2 not installed. Install with: pip install PyPDF2"
    except Exception as e:
        return f"Error reading PDF file: {e}"


def read_docx_file(file_path: str) -> str:
    """Read a Word document and extract text.
    
    Args:
        file_path: Path to the DOCX file
        
    Returns:
        Extracted text from document
    """
    try:
        from docx import Document
        
        doc = Document(file_path)
        
        result = f"Word document with {len(doc.paragraphs)} paragraphs\n\n"
        
        # Extract all text
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        
        result += "\n".join(full_text)
        
        return result
    except ImportError:
        return "python-docx not installed. Install with: pip install python-docx"
    except Exception as e:
        return f"Error reading DOCX file: {e}"


def read_pptx_file(file_path: str) -> str:
    """Read a PowerPoint presentation and extract text.
    
    Args:
        file_path: Path to the PPTX file
        
    Returns:
        Extracted text from slides
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        
        result = f"PowerPoint presentation with {len(prs.slides)} slides\n\n"
        
        for i, slide in enumerate(prs.slides, 1):
            result += f"--- Slide {i} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    result += f"{shape.text}\n"
            result += "\n"
        
        return result
    except ImportError:
        return "python-pptx not installed. Install with: pip install python-pptx"
    except Exception as e:
        return f"Error reading PPTX file: {e}"


def describe_image(file_path: str) -> str:
    """Get basic information about an image file.
    
    Args:
        file_path: Path to the image file
        
    Returns:
        Image information (dimensions, format, etc.)
    """
    try:
        from PIL import Image
        
        with Image.open(file_path) as img:
            result = f"Image file information:\n"
            result += f"Format: {img.format}\n"
            result += f"Size: {img.size[0]}x{img.size[1]} pixels\n"
            result += f"Mode: {img.mode}\n"
            
            # Get EXIF data if available
            if hasattr(img, '_getexif') and img._getexif():
                result += "\nEXIF data available"
            
            return result
    except ImportError:
        return "Pillow not installed. Install with: pip install Pillow"
    except Exception as e:
        return f"Error reading image file: {e}"


def read_audio_metadata(file_path: str) -> str:
    """Get metadata from an audio file.
    
    Args:
        file_path: Path to the audio file
        
    Returns:
        Audio file metadata
    """
    try:
        import mutagen
        
        audio = mutagen.File(file_path)
        
        result = f"Audio file information:\n"
        result += f"Length: {audio.info.length:.2f} seconds\n"
        
        if hasattr(audio.info, 'bitrate'):
            result += f"Bitrate: {audio.info.bitrate} bps\n"
        
        if hasattr(audio.info, 'sample_rate'):
            result += f"Sample rate: {audio.info.sample_rate} Hz\n"
        
        if audio.tags:
            result += "\nTags:\n"
            for key, value in audio.tags.items():
                result += f"  {key}: {value}\n"
        
        return result
    except ImportError:
        return "mutagen not installed. Install with: pip install mutagen"
    except Exception as e:
        return f"Error reading audio file: {e}"


def list_zip_contents(file_path: str) -> str:
    """List contents of a ZIP archive.
    
    Args:
        file_path: Path to the ZIP file
        
    Returns:
        List of files in the archive
    """
    try:
        import zipfile
        
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            file_list = zip_ref.namelist()
            
            result = f"ZIP archive with {len(file_list)} files:\n\n"
            
            for file_name in file_list:
                info = zip_ref.getinfo(file_name)
                result += f"  {file_name} ({info.file_size} bytes)\n"
            
            return result
    except Exception as e:
        return f"Error reading ZIP file: {e}"


def extract_zip_file(file_path: str, extract_to: Optional[str] = None) -> str:
    """Extract a ZIP archive.
    
    Args:
        file_path: Path to the ZIP file
        extract_to: Directory to extract to (default: same directory as zip)
        
    Returns:
        Status message
    """
    try:
        import zipfile
        from pathlib import Path
        
        if extract_to is None:
            extract_to = str(Path(file_path).parent / Path(file_path).stem)
        
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            zip_ref.extractall(extract_to)
            file_list = zip_ref.namelist()
        
        return f"Extracted {len(file_list)} files to: {extract_to}"
    except Exception as e:
        return f"Error extracting ZIP file: {e}"


def read_python_file(file_path: str) -> str:
    """Read a Python file and return its contents.
    
    Args:
        file_path: Path to the Python file
        
    Returns:
        Python code as string
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            code = f.read()
        
        # Count lines and basic structure
        lines = code.split('\n')
        imports = [l for l in lines if l.strip().startswith(('import ', 'from '))]
        functions = [l for l in lines if l.strip().startswith('def ')]
        classes = [l for l in lines if l.strip().startswith('class ')]
        
        result = f"Python file ({len(lines)} lines)\n"
        result += f"Imports: {len(imports)}\n"
        result += f"Functions: {len(functions)}\n"
        result += f"Classes: {len(classes)}\n\n"
        result += f"Contents:\n{code}"
        
        return result
    except Exception as e:
        return f"Error reading Python file: {e}"


def read_pdb_file(file_path: str) -> str:
    """Read a PDB (Protein Data Bank) file and extract basic information.
    
    Args:
        file_path: Path to the PDB file
        
    Returns:
        PDB file information
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            lines = f.readlines()
        
        # Extract header information
        header = [l for l in lines if l.startswith('HEADER')]
        title = [l for l in lines if l.startswith('TITLE')]
        atoms = [l for l in lines if l.startswith('ATOM')]
        
        result = "PDB file information:\n\n"
        
        if header:
            result += f"Header: {header[0][10:].strip()}\n"
        if title:
            result += f"Title: {title[0][10:].strip()}\n"
        
        result += f"\nNumber of atoms: {len(atoms)}\n"
        result += f"Total lines: {len(lines)}\n\n"
        result += "First 50 lines:\n"
        result += "".join(lines[:50])
        
        return result
    except Exception as e:
        return f"Error reading PDB file: {e}"


def get_file_info(file_path: str) -> str:
    """Get general information about any file.
    
    Args:
        file_path: Path to the file
        
    Returns:
        File information
    """
    try:
        from pathlib import Path
        import os
        
        path = Path(file_path)
        
        if not path.exists():
            return f"File not found: {file_path}"
        
        result = f"File information:\n"
        result += f"Name: {path.name}\n"
        result += f"Size: {path.stat().st_size} bytes\n"
        result += f"Extension: {path.suffix}\n"
        result += f"Absolute path: {path.absolute()}\n"
        
        return result
    except Exception as e:
        return f"Error getting file info: {e}"


def create_gaia_tools(file_paths: Optional[List[str]] = None) -> List[callable]:
    """Create a list of tools for GAIA benchmark with file access.
    
    Args:
        file_paths: Optional list of file paths that will be available
        
    Returns:
        List of tool functions
    """
    tools = [
        read_text_file,
        read_csv_file,
        read_excel_file,
        read_json_file,
        read_pdf_file,
        read_docx_file,
        read_pptx_file,
        describe_image,
        read_audio_metadata,
        list_zip_contents,
        extract_zip_file,
        read_python_file,
        read_pdb_file,
        get_file_info,
    ]
    
    return tools
